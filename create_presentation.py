#!/usr/bin/env python3
"""Generate the ChaosProbe thesis defense PowerPoint presentation.

The deck tells the *current* pre-registered study (campaigns C1/C2/C3 plus the
exploratory hotelReservation external-validity arm and the design-corrected
re-analysis). Every quoted number is taken verbatim from the thesis results
chapter (chapters/05-results.tex) and the per-hypothesis analysis drivers, and
every results figure is the same PNG the thesis embeds (thesis/figures/), so the
slides cannot drift from the document.

Slide structure (defense format, bar-first results):
  1. Title
  2. The operator's decision (motivation + gap)
  3. Research question, three layers, two knobs
  4. Pre-registration & provenance discipline
  5. ChaosProbe architecture (authored vs off-the-shelf)
  6. Three-layer measurement design
  7. Campaigns C1/C2/C3 + fault classes + hypothesis map
  8. H1  — dose-response of the east-west tail
  9. H2  — placement-dependence and the DNS intervention
 10. H3  — replication rescue under node-drain
 11. H4 & H5 — degenerate frontier + scorecard reliability
 12. Confirmatory family — the Holm capstone (central finding)
 13. Design-corrected re-analysis of the availability axis
 14. External validity — a second workload
 15. Threats to validity
 16. Conclusion, contributions & future work
 17. Questions
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ─────────────────────────────────────────────────
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_ORANGE = RGBColor(0xF3, 0x9C, 0x12)
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBD, 0xBD, 0xBD)
MID_GRAY = RGBColor(0x90, 0x90, 0xA0)
VERY_DARK = RGBColor(0x12, 0x12, 0x22)
TRANS_WHITE = RGBColor(0xF0, 0xF0, 0xF8)

CLR_CLI = RGBColor(0x34, 0x98, 0xDB)
CLR_ORCH = RGBColor(0x1A, 0xBC, 0x9C)
CLR_CHAOS = RGBColor(0xE7, 0x4C, 0x3C)
CLR_METRICS = RGBColor(0x2E, 0xCC, 0x71)
CLR_STORAGE = RGBColor(0x9B, 0x59, 0xB6)
CLR_OUTPUT = RGBColor(0xF3, 0x9C, 0x12)
CLR_INFRA = RGBColor(0x7F, 0x8C, 0x8D)

_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
FIG_DIR = os.path.join(_SCRIPT_DIR, "thesis", "figures")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Helper functions ────────────────────────────────────────────────
def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=14,
    bold=False,
    color=WHITE,
    alignment=PP_ALIGN.LEFT,
    font_name="Calibri",
):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rounded_box(
    slide,
    left,
    top,
    width,
    height,
    fill_color,
    text="",
    font_size=11,
    text_color=WHITE,
    bold=False,
    border_color=None,
    border_width=Pt(1),
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.bold = bold
        p.font.name = "Calibri"
        tf.margin_left = Pt(4)
        tf.margin_right = Pt(4)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=LIGHT_GRAY, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def add_bullet_frame(
    slide,
    left,
    top,
    width,
    height,
    items,
    font_size=13,
    color=LIGHT_GRAY,
    title=None,
    title_size=16,
    title_color=ACCENT_BLUE,
):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(title_size)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    for i, item in enumerate(items):
        if i == 0 and not title:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    return txBox


def add_table(
    slide,
    left,
    top,
    width,
    height,
    rows,
    cols,
    data,
    header_color=ACCENT_BLUE,
    cell_color=None,
    text_color=WHITE,
    header_text_color=WHITE,
    font_size=10,
):
    if cell_color is None:
        cell_color = RGBColor(0x2A, 0x2A, 0x3E)
    tbl_shape = slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tbl = tbl_shape.table

    row_height = Inches(height / rows)
    for ri in range(rows):
        tbl.rows[ri].height = row_height

    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell_text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ""
            tf = cell.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = cell_text
            run.font.size = Pt(font_size)
            run.font.name = "Calibri"
            run.font.bold = r == 0
            run.font.color.rgb = header_text_color if r == 0 else text_color
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = header_color if r == 0 else cell_color
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return tbl_shape


def slide_title(slide, title_text, subtitle_text=None):
    add_text_box(
        slide, 0.6, 0.3, 12.2, 0.7, title_text, font_size=30, bold=True, color=WHITE
    )
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.95), Inches(3), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()
    if subtitle_text:
        add_text_box(
            slide, 0.6, 1.05, 12.2, 0.5, subtitle_text, font_size=15, color=LIGHT_GRAY
        )


def add_fig(slide, left, top, width, filename, placeholder=None):
    """Embed a thesis figure preserving its aspect ratio (width-only sizing)."""
    path = os.path.join(FIG_DIR, filename)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))
    else:
        add_rounded_box(
            slide,
            left,
            top,
            width,
            3.0,
            VERY_DARK,
            placeholder or filename,
            11,
            MID_GRAY,
            False,
            border_color=MID_GRAY,
        )


def new_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    return slide


def verdict_chip(slide, left, top, text, color, width=3.0):
    add_rounded_box(
        slide,
        left,
        top,
        width,
        0.5,
        color,
        text,
        font_size=13,
        text_color=WHITE,
        bold=True,
    )


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════
slide = new_slide()

add_rounded_box(slide, 0, 0, 13.333, 2.4, VERY_DARK)
add_text_box(
    slide,
    0.8,
    0.85,
    11.7,
    1.0,
    "ChaosProbe",
    font_size=54,
    bold=True,
    color=ACCENT_BLUE,
    alignment=PP_ALIGN.CENTER,
)
add_text_box(
    slide,
    0.8,
    1.75,
    11.7,
    0.6,
    "Measuring Placement-Sensitive Resilience under Chaos",
    font_size=22,
    bold=True,
    color=WHITE,
    alignment=PP_ALIGN.CENTER,
)

add_text_box(
    slide,
    0.8,
    2.7,
    11.7,
    0.5,
    "A pre-registered, layered study in Kubernetes",
    font_size=18,
    color=LIGHT_GRAY,
    alignment=PP_ALIGN.CENTER,
)

# Three-line takeaway band.
band = [
    ("A single score is blind to placement", ACCENT_RED),
    (
        "Placement moves a real kernel mechanism — that does not reach the user",
        ACCENT_ORANGE,
    ),
    (
        "Where it does reach users is availability under node failure — predictably",
        ACCENT_GREEN,
    ),
]
y = 3.6
for txt, clr in band:
    add_rounded_box(
        slide,
        2.4,
        y,
        8.5,
        0.55,
        VERY_DARK,
        txt,
        font_size=14,
        text_color=clr,
        bold=True,
        border_color=clr,
    )
    y += 0.7

add_text_box(
    slide,
    0.8,
    6.25,
    11.7,
    0.4,
    "Yvo Hu  ·  MSc Thesis Defense  ·  University of Amsterdam",
    font_size=16,
    bold=True,
    color=WHITE,
    alignment=PP_ALIGN.CENTER,
)
add_text_box(
    slide,
    0.8,
    6.7,
    11.7,
    0.4,
    "Supervisor: [TODO]   ·   Examiner: [TODO]   ·   2026",
    font_size=12,
    color=MID_GRAY,
    alignment=PP_ALIGN.CENTER,
)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE OPERATOR'S DECISION (motivation + gap)
# ══════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(
    slide,
    "The Operator's Decision",
    "Pack the services onto few nodes, or spread them across many?",
)

add_rounded_box(slide, 0.6, 1.75, 5.9, 1.55, VERY_DARK, border_color=ACCENT_GREEN)
add_bullet_frame(
    slide,
    0.8,
    1.85,
    5.5,
    1.4,
    [
        "• Locality lowers inter-service latency — co-located calls stay",
        "   node-local (NetMARKS reports up to 37% faster responses)",
        "• Schedulers from Borg to Medea encode this as resource fit",
    ],
    font_size=12.5,
    color=LIGHT_GRAY,
    title="Pack — the locality literature",
    title_size=15,
    title_color=ACCENT_GREEN,
)

add_rounded_box(slide, 6.8, 1.75, 5.9, 1.55, VERY_DARK, border_color=ACCENT_ORANGE)
add_bullet_frame(
    slide,
    7.0,
    1.85,
    5.5,
    1.4,
    [
        "• Services in distinct failure domains survive the loss of any",
        "   one node (Medea topology spread)",
        "• Co-location invites resource contention (Bubble-Up, Quasar)",
    ],
    font_size=12.5,
    color=LIGHT_GRAY,
    title="Spread — the availability literature",
    title_size=15,
    title_color=ACCENT_ORANGE,
)

add_rounded_box(
    slide, 0.6, 3.55, 12.1, 1.25, RGBColor(0x2A, 0x2A, 0x3E), border_color=ACCENT_RED
)
add_bullet_frame(
    slide,
    0.85,
    3.65,
    11.7,
    1.1,
    [
        "• The Kubernetes scheduler scores nodes on resource fit, not fault isolation — it answers a different question",
        "• Chaos-engineering tools answer the resilience question — but collapse the outcome into a single aggregate score",
        "• Whether that score can even discriminate between placements — signal vs run-to-run noise — is unexamined",
    ],
    font_size=13,
    color=LIGHT_GRAY,
    title="What does today's tooling tell this operator?",
    title_size=15,
    title_color=ACCENT_RED,
)

add_rounded_box(slide, 0.6, 5.05, 12.1, 1.9, VERY_DARK, border_color=ACCENT_BLUE)
add_text_box(
    slide,
    0.85,
    5.15,
    11.7,
    0.45,
    "The research gap",
    font_size=16,
    bold=True,
    color=ACCENT_BLUE,
)
add_text_box(
    slide,
    0.85,
    5.6,
    11.7,
    1.25,
    "Two bodies of practice — placement and chaos engineering — have never been connected. Nobody has "
    "measured AT WHICH LAYER a placement effect appears under a given fault class, whether it reaches the "
    "user-visible outcome, and whether an aggregate score can see it at all. If placement moves a kernel "
    "mechanism but not the user outcome, a single score is the wrong instrument — and placement advice "
    "derived from it is unreliable.",
    font_size=14,
    color=TRANS_WHITE,
)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 3 — RESEARCH QUESTION, THREE LAYERS, TWO KNOBS
# ══════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(
    slide,
    "Research Question",
    "A fault-class × measurement-layer study — not a placement ranking",
)

add_rounded_box(slide, 0.6, 1.7, 12.1, 1.2, VERY_DARK, border_color=ACCENT_BLUE)
add_text_box(
    slide,
    0.9,
    1.85,
    11.5,
    1.0,
    "Under which chaos fault classes does pod placement measurably affect mechanism-level behaviour and "
    "user-visible outcomes in a Kubernetes microservice application — and when do aggregate resilience "
    "scores obscure those effects?",
    font_size=17,
    bold=True,
    color=WHITE,
)

# Three measurement layers.
add_text_box(
    slide,
    0.6,
    3.1,
    6.0,
    0.4,
    "An effect can appear at three layers — they need not agree:",
    font_size=14,
    bold=True,
    color=ACCENT_BLUE,
)
layers = [
    ("(a) Aggregate score", "one resilience number", ACCENT_RED),
    ("(b) Kernel / network mechanism", "conntrack, kube-proxy SLO", ACCENT_ORANGE),
    ("(c) User-visible outcome", "route latency + error rate", ACCENT_GREEN),
]
y = 3.6
for name, sub, clr in layers:
    add_rounded_box(slide, 0.6, y, 6.0, 0.62, VERY_DARK, border_color=clr)
    add_text_box(
        slide, 0.8, y + 0.04, 4.0, 0.5, name, font_size=13.5, bold=True, color=clr
    )
    add_text_box(slide, 4.55, y + 0.07, 2.0, 0.4, sub, font_size=11, color=MID_GRAY)
    y += 0.74

# Two knobs + arms.
add_text_box(
    slide,
    7.0,
    3.1,
    5.7,
    0.4,
    "Placement is manipulated by two orthogonal knobs:",
    font_size=14,
    bold=True,
    color=ACCENT_BLUE,
)
add_rounded_box(slide, 7.0, 3.6, 5.7, 0.78, VERY_DARK, border_color=ACCENT_PURPLE)
add_text_box(
    slide,
    7.2,
    3.66,
    5.3,
    0.7,
    "f  — cross-node dependency-edge fraction\n     (packing f=0  →  spreading f=1)",
    font_size=13,
    bold=True,
    color=WHITE,
)
add_rounded_box(slide, 7.0, 4.5, 5.7, 0.78, VERY_DARK, border_color=ACCENT_PURPLE)
add_text_box(
    slide,
    7.2,
    4.56,
    5.3,
    0.7,
    "r  — replication degree\n     (packed vs anti-affine)",
    font_size=13,
    bold=True,
    color=WHITE,
)
add_text_box(
    slide,
    7.0,
    5.4,
    5.7,
    0.8,
    "+ two interventional arms: NodeLocal DNSCache (executed) and kube-proxy mode (de-scoped to ipvs, not run).",
    font_size=12,
    color=LIGHT_GRAY,
)

add_rounded_box(
    slide, 0.6, 6.35, 12.1, 0.75, RGBColor(0x2A, 0x2A, 0x3E), border_color=ACCENT_GREEN
)
add_text_box(
    slide,
    0.85,
    6.45,
    11.7,
    0.6,
    "Reliability first (H5): if a single score could rank placements reliably, the layered design would be a "
    "refinement. Because its availability layer cannot in this regime, the layered design is what produces the findings.",
    font_size=12.5,
    color=TRANS_WHITE,
)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 — PRE-REGISTRATION & PROVENANCE DISCIPLINE
# ══════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(
    slide,
    "Pre-Registration & Provenance",
    "The methodological backbone — removing the freedom to choose tests after seeing the data",
)

add_rounded_box(slide, 0.6, 1.7, 5.9, 4.6, VERY_DARK, border_color=ACCENT_BLUE)
add_bullet_frame(
    slide,
    0.85,
    1.85,
    5.5,
    4.4,
    [
        "• Hypotheses, smallest effect sizes of interest (SESOIs),",
        "   per-cell sample sizes, and the analysis CODE were frozen",
        "   and deposited BEFORE any confirmatory data was collected",
        "",
        "• Frozen pre-registration deposited at its freeze commit",
        "   (DOI 10.5281/zenodo.20690836, commit 20097c1)",
        "",
        "• Each member carries a registered effect-size / reliability",
        "   bar — Holm-significance is necessary, not sufficient",
        "",
        "• The one non-blind deviation (D3) is logged in full in the",
        "   deviations log — disclosed, not hidden",
    ],
    font_size=13,
    color=LIGHT_GRAY,
    title="Frozen before collection",
    title_size=16,
    title_color=ACCENT_BLUE,
)

add_rounded_box(slide, 6.8, 1.7, 5.9, 4.6, VERY_DARK, border_color=ACCENT_GREEN)
add_bullet_frame(
    slide,
    7.05,
    1.85,
    5.5,
    4.4,
    [
        "• Independent single-commit sessions are the unit of",
        "   analysis — no pooling across code versions",
        "",
        "• Every session gated by doctor --strict: clean tree,",
        "   scenario-hash present, runMetadata present, git.dirty=false",
        "",
        "• Deposit-before-analysis: each campaign's raw data is",
        "   deposited under a DOI before its analysis is written",
        "",
        "• Every quoted number traces to an archived, hash-stamped",
        "   run (artifact-manifest.json, SHA-256 per file)",
    ],
    font_size=13,
    color=LIGHT_GRAY,
    title="Provenance-gated replication",
    title_size=16,
    title_color=ACCENT_GREEN,
)

add_rounded_box(
    slide, 0.6, 6.45, 12.1, 0.65, RGBColor(0x2A, 0x2A, 0x3E), border_color=ACCENT_ORANGE
)
add_text_box(
    slide,
    0.85,
    6.53,
    11.7,
    0.5,
    "This converts “we observed placement effects” into the more honest “here is exactly which effects "
    "survive a strict bar, and which do not.”",
    font_size=13,
    bold=True,
    color=TRANS_WHITE,
)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 — CHAOSPROBE ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(
    slide,
    "ChaosProbe — The Framework",
    "Contribution 1: a placement-aware, pre-registered, provenance-gated chaos-evaluation method",
)

authored = [
    (
        "Placement solver",
        "drives the two knobs (f, r) by\ndeployment-level nodeSelector mutation",
        CLR_ORCH,
    ),
    (
        "Cross-layer probers",
        "Prometheus mechanism metrics,\nroute latency (dependent vs control),\nEndpointSlice availability snapshots",
        CLR_METRICS,
    ),
    (
        "Cross-node fraction",
        "computed from the Neo4j dependency\ngraph + the realized placement",
        CLR_STORAGE,
    ),
    (
        "Layered scorecard",
        "availability / mechanism / user-tail\nsub-scores — not one number",
        CLR_OUTPUT,
    ),
]
x = 0.6
for name, sub, clr in authored:
    add_rounded_box(slide, x, 1.8, 2.92, 1.7, VERY_DARK, border_color=clr)
    add_text_box(
        slide, x + 0.15, 1.9, 2.62, 0.5, name, font_size=14, bold=True, color=clr
    )
    add_text_box(slide, x + 0.15, 2.4, 2.62, 1.05, sub, font_size=11, color=LIGHT_GRAY)
    x += 3.04

add_text_box(
    slide,
    0.6,
    3.75,
    12.1,
    0.4,
    "↓   driving and reading off-the-shelf infrastructure   ↓",
    font_size=13,
    bold=True,
    color=MID_GRAY,
    alignment=PP_ALIGN.CENTER,
)

infra = [
    ("LitmusChaos\n+ ChaosCenter", "fault execution"),
    ("Prometheus\n+ Locust", "metrics & load"),
    ("Neo4j 5", "graph storage"),
    ("Online Boutique", "system under test"),
]
x = 0.6
for name, sub in infra:
    add_rounded_box(
        slide, x, 4.25, 2.92, 1.0, RGBColor(0x24, 0x24, 0x36), border_color=CLR_INFRA
    )
    add_text_box(
        slide, x + 0.15, 4.33, 2.62, 0.55, name, font_size=12.5, bold=True, color=WHITE
    )
    add_text_box(slide, x + 0.15, 4.88, 2.62, 0.35, sub, font_size=10.5, color=MID_GRAY)
    x += 3.04

add_rounded_box(slide, 0.6, 5.55, 12.1, 1.45, VERY_DARK, border_color=ACCENT_BLUE)
add_text_box(
    slide,
    0.85,
    5.65,
    11.7,
    0.4,
    "What the candidate did vs what the tooling did",
    font_size=15,
    bold=True,
    color=ACCENT_BLUE,
)
add_text_box(
    slide,
    0.85,
    6.05,
    11.7,
    0.95,
    "Authored: the placement solver and two-knob design; the three-layer measurement design with its "
    "dependent-vs-control route confound check; the cross-node-fraction metric and EndpointSlice trough; the "
    "layered scorecard; the pre-registration, doctor --strict gate, discard-not-patch rule, deposit-before-"
    "analysis protocol; and every analysis script.   Integrated off-the-shelf: LitmusChaos, Prometheus/Locust, "
    "Neo4j, and the Online Boutique application execute the science — they do not define it.",
    font_size=12.5,
    color=TRANS_WHITE,
)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 6 — THREE-LAYER MEASUREMENT DESIGN
# ══════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(
    slide,
    "Three-Layer Measurement Design",
    "60s steady state  →  120s fault active  →  60s recovery — read at every layer",
)

rows = [
    ["Layer", "What is measured", "How"],
    [
        "Aggregate score",
        "Layered resilience scorecard (mean of probe success %), reported with bootstrap CIs",
        "LitmusChaos resilience probes + 5 Rust cmdProbes",
    ],
    [
        "Kernel / network\nmechanism",
        "UDP-conntrack drop, kube-proxy SLO p99, CoreDNS p99, CPU throttle (SIG-Scalability)",
        "Prometheus PromQL, churn-prober",
    ],
    [
        "User-visible\noutcome",
        "HTTP route p95 latency + error rate, split dependent-vs-control route; east-west p95",
        "kubectl exec → wget/python3, 3.5s interval",
    ],
    [
        "Availability\n(node-drain)",
        "EndpointSlice trough = fraction of app-ready endpoints lost when the node drains",
        "K8s EndpointSlice snapshots",
    ],
]
add_table(
    slide,
    0.6,
    1.75,
    8.1,
    3.5,
    len(rows),
    3,
    rows,
    header_color=ACCENT_BLUE,
    font_size=11.5,
)

add_rounded_box(slide, 8.95, 1.75, 3.75, 3.5, VERY_DARK, border_color=ACCENT_PURPLE)
add_bullet_frame(
    slide,
    9.15,
    1.85,
    3.4,
    3.35,
    [
        "• 6 continuous probers run as background threads",
        "   through all three phases",
        "",
        "• Recovery split: deletion→scheduled (d2s) vs",
        "   scheduled→ready (s2r)",
        "",
        "• Cross-node fraction f computed pre-chaos",
        "   from the dependency graph + placement",
        "",
        "• The layered scorecard separates the layer",
        "   that carries signal from the one that does not",
        "   — rather than averaging them into one number",
    ],
    font_size=11.5,
    color=LIGHT_GRAY,
    title="The instrument",
    title_size=14,
    title_color=ACCENT_PURPLE,
)

add_rounded_box(
    slide, 0.6, 5.5, 12.1, 1.5, RGBColor(0x2A, 0x2A, 0x3E), border_color=ACCENT_GREEN
)
add_text_box(
    slide,
    0.85,
    5.6,
    11.7,
    0.4,
    "Why three layers, not one score",
    font_size=15,
    bold=True,
    color=ACCENT_GREEN,
)
add_text_box(
    slide,
    0.85,
    6.0,
    11.7,
    1.0,
    "A single aggregate score's failure modes compound rather than cancel. Under churn it is NOISY — the "
    "availability layer it leans on has little reproducible between-condition signal. Under node-drain it is "
    "ABSENT — the drain takes down the infrastructure the probes depend on, every probe returns Unknown, and "
    "the score is undefined exactly when the availability outcome is most dramatic. So availability is read "
    "directly from EndpointSlice troughs.",
    font_size=12.5,
    color=TRANS_WHITE,
)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 7 — CAMPAIGNS + FAULT CLASSES + HYPOTHESIS MAP
# ══════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(
    slide,
    "Campaigns & Hypotheses",
    "Three pre-registered campaigns, each deposited under a DOI before analysis",
)

camp = [
    ["Campaign", "Fault class", "Design", "Sess.", "Tests", "DOI"],
    [
        "C1 dose-response",
        "churn (pod-delete)",
        "complete-block f∈{0,.25,.5,.75,1}, r=1, 5×3 iters",
        "8",
        "H1, H4, H5",
        "20690737",
    ],
    [
        "C2 replication rescue",
        "node failure (drain)",
        "r∈{1,3} × mode (packed/anti-affine), 3 cells",
        "24",
        "H3",
        "20726729",
    ],
    [
        "C3 placement + DNS",
        "churn (pod-delete)",
        "f∈{0,1} × dnsCache off/on, 7 matched pairs",
        "14",
        "H2",
        "20748970",
    ],
]
add_table(
    slide,
    0.6,
    1.75,
    12.1,
    2.0,
    len(camp),
    6,
    camp,
    header_color=ACCENT_BLUE,
    font_size=11,
)

# Hypothesis bar-first map.
hyps = [
    ["H", "Claim", "Registered bar"],
    [
        "H1",
        "Dose-response of the east-west tail",
        "monotone increase, ≥ 15% SESOI (f=0→1)",
    ],
    [
        "H2",
        "Placement-dependence + DNS intervention",
        "spread>packed drop AND DNS shrinks ≥50%",
    ],
    [
        "H3",
        "Replication rescue under node-drain",
        "r×mode interaction + anti-affine rescue margin",
    ],
    [
        "H4",
        "Placement frontier (descriptive)",
        "non-dominated set under frozen margins",
    ],
    ["H5", "Layered scorecard reliability", "availability AND mechanism ICC ≥ 0.5"],
    [
        "H6",
        "Direction transfer (iptables mode)",
        "not attempted — cluster ran ipvs, de-scoped",
    ],
]
add_table(
    slide,
    0.6,
    4.0,
    9.4,
    3.0,
    len(hyps),
    3,
    hyps,
    header_color=ACCENT_PURPLE,
    font_size=11,
)

add_rounded_box(slide, 10.2, 4.0, 2.5, 3.0, VERY_DARK, border_color=ACCENT_ORANGE)
add_bullet_frame(
    slide,
    10.35,
    4.1,
    2.2,
    2.85,
    [
        "Two fault classes:",
        "",
        "• CHURN (pod-delete)",
        "   single replica gone",
        "",
        "• NODE FAILURE (drain)",
        "   placement-dependent",
        "   blast radius",
        "",
        "Cluster: 8×2-vCPU/",
        "4-GiB workers,",
        "K8s v1.28.6, ipvs",
    ],
    font_size=11,
    color=LIGHT_GRAY,
    title="Faults & cluster",
    title_size=13,
    title_color=ACCENT_ORANGE,
)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 8 — RESULTS: H1
# ══════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(
    slide,
    "H1 — Dose-Response of the East-West Tail",
    "Bar: east-west p95 rises monotonically in f, by ≥ 15% (f=0→f=1)",
)

add_fig(slide, 6.7, 1.7, 6.4, "fig-h1-dose-response.png")

add_bullet_frame(
    slide,
    0.6,
    1.8,
    6.0,
    2.5,
    [
        "• Page's L trend test:  L = 410,  z = 3.54,  p = 0.0002",
        "• Per-level median p95 (ms):",
        "     35.74 / 38.60 / 41.40 / 39.55 / 40.51",
        "• f=0 → f=1 effect:  +13.35%",
        "• Medians rise, then plateau/dip at f=0.75",
    ],
    font_size=14,
    color=LIGHT_GRAY,
    title="Result — trend detected, but below the SESOI",
    title_size=16,
    title_color=ACCENT_BLUE,
)

verdict_chip(
    slide, 0.6, 4.35, "H1 NOT SUPPORTED — below the 15% SESOI", ACCENT_RED, width=6.0
)

add_rounded_box(slide, 0.6, 5.05, 6.0, 1.95, VERY_DARK, border_color=ACCENT_ORANGE)
add_text_box(
    slide,
    0.8,
    5.15,
    5.6,
    0.4,
    "Disclosed deviation (D3, non-blind)",
    font_size=13,
    bold=True,
    color=ACCENT_ORANGE,
)
add_text_box(
    slide,
    0.8,
    5.55,
    5.6,
    1.4,
    "The registered pre-window UDP-slope taint flags every iteration at f=0.25 and f=0.5 — zero complete "
    "blocks, so Page's L is uncomputable and H1 is unrunnable as registered. The taint is withdrawn on "
    "protocol-independent grounds (it gates DNS/UDP, a different protocol from the TCP/gRPC outcome). This is "
    "the study's one outcome-aware deviation — logged in full.",
    font_size=11,
    color=TRANS_WHITE,
)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 9 — RESULTS: H2
# ══════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(
    slide,
    "H2 — Placement-Dependence and the DNS Intervention",
    "Bar (both must pass): (a) spread > packed conntrack drop  AND  (b) DNS cache shrinks it ≥ 50%",
)

add_fig(slide, 0.6, 1.75, 7.4, "fig-h2-conntrack-dns.png")

add_bullet_frame(
    slide,
    8.2,
    1.8,
    4.7,
    3.4,
    [
        "(a) Placement — REVERSED:",
        "   packed > spread in 7/7 pairs",
        "   (packed 11153, spread 8929; p=0.99)",
        "",
        "(b) DNS mechanism — PASSES:",
        "   78.0% median shrink (all 7 pairs 61–85%)",
        "   packed drop also → ≈0 with cache on",
    ],
    font_size=13.5,
    color=LIGHT_GRAY,
    title="Result",
    title_size=16,
    title_color=ACCENT_BLUE,
)

verdict_chip(
    slide, 8.2, 4.85, "CONJUNCTION = FALSE — H2 not supported", ACCENT_RED, width=4.7
)

add_rounded_box(
    slide, 0.6, 5.5, 12.3, 1.5, RGBColor(0x2A, 0x2A, 0x3E), border_color=ACCENT_GREEN
)
add_text_box(
    slide,
    0.85,
    5.6,
    12.0,
    0.4,
    "The reversal sharpens the mechanism",
    font_size=15,
    bold=True,
    color=ACCENT_GREEN,
)
add_text_box(
    slide,
    0.85,
    6.0,
    12.0,
    1.0,
    "Co-locating a service's replicas on one node CONCENTRATES the per-node conntrack churn when pod-delete "
    "hits that node — a larger drop than spreading the same replicas. The DNS cache removes the drop under "
    "BOTH placements, so the floodable pool is the general per-node DNS/UDP conntrack churn, not specifically "
    "cross-node DNS. Placement does matter for conntrack — significantly, but with the opposite sign to the "
    "prediction. The DNS cache is the one actionable lever the study confirmed.",
    font_size=12.5,
    color=TRANS_WHITE,
)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 — RESULTS: H3
# ══════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(
    slide,
    "H3 — Replication Rescue under Node-Drain",
    "Bar: r×mode interaction + anti-affine rescue clears its margin + packed r=3 ≈ r=1",
)

add_fig(slide, 0.6, 1.75, 7.4, "fig-h3-replication-rescue.png")

rows = [
    ["co-primary", "r1", "r3-pk", "r3-anti", "vs bar"],
    [
        "trough depth",
        "0.091",
        "0.091",
        "0.046",
        "interaction p=.0065; rescue < .091 margin → not met",
    ],
    [
        "user-route error",
        "0.632",
        "0.632",
        "0.000",
        "p≈0; rescue ≥ .302 → MET; TOST in band",
    ],
]
add_table(
    slide,
    8.1,
    1.8,
    4.8,
    1.7,
    len(rows),
    5,
    rows,
    header_color=ACCENT_BLUE,
    font_size=10,
)

verdict_chip(
    slide,
    8.1,
    3.7,
    "CONJUNCTION = FALSE — but by construction",
    ACCENT_ORANGE,
    width=4.8,
)

add_rounded_box(slide, 8.1, 4.35, 4.8, 2.65, VERY_DARK, border_color=ACCENT_GREEN)
add_text_box(
    slide,
    8.3,
    4.45,
    4.4,
    0.4,
    "The error face rescues",
    font_size=14,
    bold=True,
    color=ACCENT_GREEN,
)
add_text_box(
    slide,
    8.3,
    4.85,
    4.4,
    2.1,
    "Anti-affine r=3 fully rescues the user-route error (0.632 → 0.0) — strong, significant directional "
    "support. The depth face cannot adjudicate: under round-robin spread, draining one node costs r=1 only "
    "≈1 pod, so the r=1 depth (0.091) EQUALS the one-pod margin it must beat. The metric's dynamic range "
    "coincides with the bar — a measurement-design limit surfaced honestly, reported exactly as registered "
    "(not retuned). The design-fix slide shows the rescue is real once this artifact is removed.",
    font_size=11,
    color=TRANS_WHITE,
)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 11 — RESULTS: H4 & H5
# ══════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(
    slide,
    "H4 & H5 — Frontier and Scorecard Reliability",
    "The availability axis is dead under pod-delete — and that is diagnostic, not a method failure",
)

# H4
add_fig(slide, 0.6, 1.8, 4.3, "fig-h4-frontier.png")
add_bullet_frame(
    slide,
    0.6,
    5.05,
    5.9,
    1.95,
    [
        "• All 5/5 placements non-dominated",
        "• Latency face varies (35.7–41.4 ms — the H1 signal)",
        "• Availability face DEGENERATE: pod-delete removes",
        "   one pod, so trough depth ≈ 1 pod for every placement",
        "• The trade-off the frontier was designed to reveal is",
        "   not realizable from this data",
    ],
    font_size=12,
    color=LIGHT_GRAY,
    title="H4 — frontier (descriptive): degenerate",
    title_size=14,
    title_color=ACCENT_ORANGE,
)

# H5
add_text_box(
    slide,
    7.0,
    1.55,
    5.9,
    0.35,
    "H5 — layered scorecard reliability",
    font_size=14,
    bold=True,
    color=ACCENT_BLUE,
)
add_fig(slide, 8.05, 1.95, 4.6, "fig-h5-scorecard-icc.png")
rows = [
    ["sub-score", "ICC", "≥0.5?"],
    ["availability (required)", "0.180", "no"],
    ["mechanism (required)", "0.994", "yes"],
    ["naive aggregate baseline", "0.066", "—"],
]
add_table(
    slide,
    7.0,
    5.0,
    5.9,
    1.25,
    len(rows),
    3,
    rows,
    header_color=ACCENT_BLUE,
    font_size=10.5,
)
add_text_box(
    slide,
    7.0,
    6.35,
    5.9,
    0.65,
    "Required conjunction FAILS (availability ICC < 0.5). But the mechanism layer is far more "
    "reliable than a single aggregate (0.994 vs 0.066) — the headline reliability result.",
    font_size=11,
    color=TRANS_WHITE,
)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 12 — FAMILY HOLM CAPSTONE (CENTRAL FINDING)
# ══════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(
    slide,
    "The Confirmatory Family — Holm Capstone",
    "The four confirmatory hypotheses corrected together at α = 0.05",
)

rows = [
    ["hyp", "primary test", "input p", "Holm-adj p", "supported?"],
    ["H1", "dose-response (Page's L)", "0.0002", "0.0008", "no — effect < SESOI"],
    ["H2", "placement + DNS", "0.98875", "0.98875", "no — placement reversed"],
    ["H3", "replication rescue", "0.0065", "0.0195", "no — rescue margin unmet"],
    ["H5", "layered scorecard ICC", "0.2501", "0.5002", "no — availability ICC < 0.5"],
]
add_table(
    slide,
    0.6,
    1.75,
    8.2,
    2.6,
    len(rows),
    5,
    rows,
    header_color=ACCENT_PURPLE,
    font_size=11,
)

add_rounded_box(slide, 9.0, 1.75, 3.9, 2.6, VERY_DARK, border_color=ACCENT_RED)
add_text_box(
    slide,
    9.2,
    1.9,
    3.5,
    2.4,
    "Family verdict:\n\nNo confirmatory hypothesis is supported.\n\nTwo primaries (H1's trend, H3's "
    "interaction) survive Holm — but Holm-significance is necessary, not sufficient. Every member also "
    "fails its registered effect-size / reliability bar.",
    font_size=12.5,
    bold=True,
    color=WHITE,
)

add_rounded_box(
    slide, 0.6, 4.55, 12.3, 2.45, RGBColor(0x2A, 0x2A, 0x3E), border_color=ACCENT_GREEN
)
add_text_box(
    slide,
    0.85,
    4.65,
    12.0,
    0.45,
    "The central finding: a real mechanism, rigorously bounded in reach",
    font_size=16,
    bold=True,
    color=ACCENT_GREEN,
)
add_text_box(
    slide,
    0.85,
    5.15,
    12.0,
    1.8,
    "The consistent positive thread across all three campaigns is the MECHANISM: the UDP-conntrack "
    "reconvergence signature is real and individually significant every time it is measured — placement-"
    "dependent under churn (H2), interacting with replication under node-drain (H3's error face), and the one "
    "scorecard layer with near-perfect test-retest reliability (H5, mechanism ICC 0.994). What does NOT hold, "
    "at the pre-registered bar, is that this mechanism translates into the user-visible placement, "
    "availability, and dose-response advantages the hypotheses predicted. Once pod-delete at r=1 is understood "
    "as a churn fault, the bound is the expected shape: during the kill window the single replica is simply "
    "gone, and every dependent request fails identically whether services are packed or spread.",
    font_size=13,
    color=TRANS_WHITE,
)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 13 — DESIGN-CORRECTED RE-ANALYSIS
# ══════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(
    slide,
    "Design-Corrected Re-Analysis (Exploratory)",
    "The availability-axis tests were construction-limited — node-drain makes the axis live",
)

add_fig(slide, 0.6, 1.75, 7.4, "fig-design-fix-availability.png")

add_rounded_box(slide, 8.2, 1.75, 4.7, 1.5, VERY_DARK, border_color=ACCENT_ORANGE)
add_text_box(
    slide,
    8.4,
    1.83,
    4.3,
    1.35,
    "The fault, not the finding: pod-delete at r=1 removes the only replica — trough ≈ 1 pod for EVERY "
    "placement. That made H3's depth un-passable, H4's frontier degenerate, and H5's availability ICC a "
    "measurement of noise.",
    font_size=11.5,
    color=TRANS_WHITE,
)

add_bullet_frame(
    slide,
    8.2,
    3.4,
    4.7,
    2.55,
    [
        "• C4: node-drain dose-response, 8 sessions,",
        "   criteria pre-declared before data examined",
        "• Availability trough now 1.00 (packed) → 0.36",
        "   (spread) — a LIVE axis, a real trade-off",
        "• H3: user error 0.63 → 0 + significant interaction",
        "   (p=0.0065); depth halves but knife-edge — the",
        "   robust evidence is the interaction, not the threshold",
        "• H5 availability ICC under drain is 1.0 BY",
        "   CONSTRUCTION — large signal, not a retest estimate",
    ],
    font_size=11.5,
    color=LIGHT_GRAY,
    title="What the correction shows",
    title_size=14,
    title_color=ACCENT_GREEN,
)

add_rounded_box(
    slide, 0.6, 6.1, 12.3, 0.9, RGBColor(0x2A, 0x2A, 0x3E), border_color=ACCENT_BLUE
)
add_text_box(
    slide,
    0.85,
    6.18,
    12.0,
    0.75,
    "Honest scope: this does NOT re-open the frozen confirmatory verdicts. It shows H3's not-supported and "
    "H4's degenerate results were partly construction artifacts whose underlying effects are real once "
    "availability can move — exploratory, outside the Holm family, deposited to the same bar.",
    font_size=12,
    bold=True,
    color=TRANS_WHITE,
)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 14 — EXTERNAL VALIDITY
# ══════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(
    slide,
    "External Validity — A Second Workload",
    "Exploratory replication on DeathStarBench hotelReservation (deep gRPC/Consul fan-out)",
)

add_fig(slide, 0.6, 1.75, 7.4, "fig-hotel-external-validity.png")

rows = [
    ["study", "registered bar", "hotel outcome"],
    ["H1 dose-response", "monotone, ≥15% SESOI", "Page's L p=0.99 — NO increase"],
    ["H3 repl. rescue", "interaction + margins", "CONJUNCTION = False"],
]
add_table(
    slide,
    8.2,
    1.8,
    4.7,
    1.5,
    len(rows),
    3,
    rows,
    header_color=ACCENT_BLUE,
    font_size=10.5,
)

add_bullet_frame(
    slide,
    8.2,
    3.5,
    4.7,
    1.95,
    [
        "• 32 sessions, all doctor --strict clean",
        "• H1: no dose-response (tight 5–9 ms band,",
        "   if anything a mild decrease)",
        "• H3: significant r×mode interaction, anti-affine",
        "   directionally best — but here NEITHER margin is",
        "   cleared (depth 0.044, error 0.212), a broader",
        "   miss than online-boutique's clean error rescue",
    ],
    font_size=11.5,
    color=LIGHT_GRAY,
    title="Both arms corroborate online-boutique",
    title_size=14,
    title_color=ACCENT_GREEN,
)

add_rounded_box(
    slide, 0.6, 6.1, 12.3, 0.9, RGBColor(0x2A, 0x2A, 0x3E), border_color=ACCENT_GREEN
)
add_text_box(
    slide,
    0.85,
    6.18,
    12.0,
    0.75,
    "A second, structurally different workload reproduces both the below-SESOI dose-response and the "
    "conjunction-false rescue verdict: the mechanism (spreading shrinks the blast radius, speeds recovery) is "
    "robust on both applications; the strong placement and margin-clearing claims hold on neither. The "
    "strongest available evidence the finding is not an artifact of one topology. (DOI 10.5281/zenodo.20792129)",
    font_size=12,
    color=TRANS_WHITE,
)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 15 — THREATS TO VALIDITY
# ══════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(
    slide, "Threats to Validity", "The boundary of every claim — stated, not hidden"
)

add_rounded_box(slide, 0.6, 1.75, 5.95, 5.25, VERY_DARK, border_color=ACCENT_ORANGE)
add_bullet_frame(
    slide,
    0.85,
    1.9,
    5.5,
    5.05,
    [
        "• Single-replica baseline — pod-delete at r=1 is a pure",
        "   churn fault; production runs multiple replicas. The",
        "   design-fix (node-drain) addresses the availability axis",
        "",
        "• One workload — partly addressed: hotelReservation",
        "   replicates both placement-bearing arms",
        "",
        "• Virtualized cluster — Vagrant/libvirt KVM; bare-metal",
        "   I/O may differ. Conntrack is K8s-version-pinned (v1.28.6)",
        "",
        "• IPVS only — the iptables-mode direction-transfer",
        "   comparison (H6) was de-scoped before collection",
        "",
        "• One non-blind deviation (D3) — outcome-aware taint",
        "   withdrawal; disclosed in full, sensitivity check available",
    ],
    font_size=12.5,
    color=LIGHT_GRAY,
    title="What bounds the claims",
    title_size=15,
    title_color=ACCENT_ORANGE,
)

add_rounded_box(slide, 6.85, 1.75, 5.85, 5.25, VERY_DARK, border_color=ACCENT_GREEN)
add_bullet_frame(
    slide,
    7.1,
    1.9,
    5.4,
    5.05,
    [
        "• Pre-registration fixes the SESOIs against a measured",
        "   noise floor BEFORE collection — no test-after-seeing-data",
        "",
        "• Holm correction across the confirmatory family",
        "",
        "• doctor --strict gate on every session; discard-not-patch",
        "",
        "• Deposit-before-analysis under DOIs; every number traces",
        "   to a hash-stamped archived run",
        "",
        "• Dependent-vs-control route split controls the user layer",
        "   for run-level confounds",
        "",
        "• A second workload corroborates the central finding",
    ],
    font_size=12.5,
    color=LIGHT_GRAY,
    title="What protects the claims",
    title_size=15,
    title_color=ACCENT_GREEN,
)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 16 — CONCLUSION, CONTRIBUTIONS & FUTURE WORK
# ══════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(
    slide,
    "Conclusion & Future Work",
    "Chaos evaluation of placement needs layered measurement, pre-registration, and provenance — not a single score",
)

add_rounded_box(slide, 0.6, 1.7, 6.0, 2.55, VERY_DARK, border_color=ACCENT_BLUE)
add_bullet_frame(
    slide,
    0.85,
    1.8,
    5.6,
    2.4,
    [
        "1. A pre-registered, provenance-gated measurement",
        "    methodology (ChaosProbe) — reusable, outcome-",
        "    independent",
        "2. A positively-established, precisely-bounded mechanism",
        "    — real conntrack signature, DNS cache removes it,",
        "    does not reach the user at the registered bar",
        "3. A reproducible campaign protocol with DOI-deposited",
        "    artifacts (deposit-before-analysis)",
    ],
    font_size=12.5,
    color=LIGHT_GRAY,
    title="Three contributions",
    title_size=15,
    title_color=ACCENT_BLUE,
)

add_rounded_box(slide, 6.8, 1.7, 5.9, 2.55, VERY_DARK, border_color=ACCENT_GREEN)
add_bullet_frame(
    slide,
    7.05,
    1.8,
    5.5,
    2.4,
    [
        "• H1: aggregate score blind; dose-response below SESOI",
        "• H2: placement moves conntrack (reversed sign); DNS",
        "   cache is the actionable lever",
        "• H3: anti-affine rescues the user-error face under drain",
        "• H5: mechanism ICC 0.994 vs naive 0.066 — audit a",
        "   score's reliability before trusting it to rank",
        "• Design-fix: availability bites users under node failure",
        "   — predictable from the dependency graph",
    ],
    font_size=12,
    color=LIGHT_GRAY,
    title="Key findings",
    title_size=15,
    title_color=ACCENT_GREEN,
)

add_rounded_box(slide, 0.6, 4.45, 12.1, 1.45, VERY_DARK, border_color=ACCENT_PURPLE)
add_bullet_frame(
    slide,
    0.85,
    4.55,
    11.7,
    1.3,
    [
        "• Replicate the family on different INFRASTRUCTURE (managed cluster, iptables/nftables proxy, bare metal)",
        "• A rescue margin the depth face can express (relative to realized r=1 depth, or integrated depth×duration)",
        "• Apportion the conntrack drop: kernel TCP teardown vs kube-proxy's UDP-only cleanup (composition probe)",
        "• Integrate the cross-node fraction into a scheduler as a scoring plugin",
    ],
    font_size=12,
    color=LIGHT_GRAY,
    title="Future work",
    title_size=15,
    title_color=ACCENT_PURPLE,
)

add_rounded_box(
    slide, 0.6, 6.1, 12.1, 0.9, RGBColor(0x2A, 0x2A, 0x3E), border_color=ACCENT_GREEN
)
add_text_box(
    slide,
    0.85,
    6.2,
    11.7,
    0.75,
    "A single score is blind to placement (H1). Placement acts at the mechanism layer (H2) without reaching "
    "the user (H3); where it does reach users is availability under node failure — predictably, from the "
    "dependency graph (design-fix). A real mechanism, rigorously bounded in reach.",
    font_size=13,
    bold=True,
    color=TRANS_WHITE,
)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 17 — QUESTIONS
# ══════════════════════════════════════════════════════════════════════
slide = new_slide()

add_text_box(
    slide,
    0.6,
    2.2,
    12.1,
    1.0,
    "Thank You",
    font_size=52,
    bold=True,
    color=ACCENT_BLUE,
    alignment=PP_ALIGN.CENTER,
)
add_text_box(
    slide,
    0.6,
    3.3,
    12.1,
    0.6,
    "Questions?",
    font_size=26,
    color=LIGHT_GRAY,
    alignment=PP_ALIGN.CENTER,
)

stats = [
    ("3", "pre-registered\ncampaigns", ACCENT_BLUE),
    ("46", "doctor-strict\nsessions", ACCENT_GREEN),
    ("0", "hypotheses supported\n(mechanism real)", ACCENT_ORANGE),
    ("0.994", "mechanism ICC\nvs 0.066 naive", ACCENT_PURPLE),
]
x = 1.5
for val, label, clr in stats:
    add_rounded_box(slide, x, 4.6, 2.4, 1.4, VERY_DARK, border_color=clr)
    add_text_box(
        slide,
        x,
        4.75,
        2.4,
        0.6,
        val,
        font_size=30,
        bold=True,
        color=clr,
        alignment=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        x,
        5.45,
        2.4,
        0.5,
        label,
        font_size=11.5,
        color=LIGHT_GRAY,
        alignment=PP_ALIGN.CENTER,
    )
    x += 2.6

add_text_box(
    slide,
    0.6,
    6.5,
    12.1,
    0.4,
    "Raw data, pre-registration, and analysis code deposited under DOIs · every number traces to an archived run",
    font_size=12,
    color=MID_GRAY,
    alignment=PP_ALIGN.CENTER,
)


# ══════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════
output_path = os.path.join(_SCRIPT_DIR, "ChaosProbe_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Slides: {len(prs.slides._sldIdLst)}")
